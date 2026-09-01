import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_pptx="CSE24562_Malayalam.pptx", screenshot_path="execution_screenshot.png"):
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    BG_COLOR = RGBColor(15, 23, 42)      # Dark slate #0f172a
    CARD_BG = RGBColor(30, 41, 59)       # Dark blue-gray #1e293b
    CARD_BORDER = RGBColor(51, 65, 85)   # Border #334155
    ACCENT_BLUE = RGBColor(56, 189, 248) # Cyan #38bdf8
    ACCENT_GREEN = RGBColor(74, 222, 128)# Green #4ade80
    TEXT_WHITE = RGBColor(248, 250, 252) # Slate white #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)# Slate muted #94a3b8

    blank_slide_layout = prs.slide_layouts[6]
    
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category_text="SLM NLP ASSIGNMENT"):
        # Header Container
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 1: How SLM Works
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)
    add_header(slide1, "How Small Language Model (SLM) Works", "SLM FUNDAMENTALS & ARCHITECTURE")
    
    # Left Card: Definition & Pipeline
    left_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG
    left_card.line.color.rgb = CARD_BORDER
    
    tf = left_card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "What is a Small Language Model?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    bullets = [
        "A Small Language Model (SLM) is a compact, highly efficient AI neural network designed for natural language understanding and generation.",
        "Typically contains under 10 Billion parameters (e.g. 500M to 3.8B params).",
        "Trained using curated, high-quality domain-specific datasets and synthetic data.",
        "Basic Pipeline: Input Text → Tokenization → Embedding & Transformer Attention → Probability Logits → Output Token Generation.",
        "Learns statistical word, sentence, and semantic relationships without massive compute requirements."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

    # Right Card: Advantages & Edge Capability
    right_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = CARD_BORDER
    
    tf2 = right_card.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = Inches(0.3)
    
    p = tf2.paragraphs[0]
    p.text = "Key Advantages of SLMs"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    adv = [
        ("Ultra-Low Latency & High Speed", "Generates responses in milliseconds due to small parameter footprint."),
        ("Low Memory & Hardware Requirements", "Can run locally on CPUs, edge devices, laptops, and mobile phones without high-end GPUs."),
        ("Cost-Efficient", "Drastically reduces cloud API inference costs and carbon footprint."),
        ("Enhanced Privacy & Data Security", "Processes sensitive technical data locally offline without sending information to cloud servers."),
        ("Easy Domain Adaptability", "Quick to fine-tune on domain-specific datasets (e.g. Malayalam translation).")
    ]
    for title, desc in adv:
        p = tf2.add_paragraph()
        p.text = f"✔ {title}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(6)
        
        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_after = Pt(4)

    # ----------------------------------------------------
    # SLIDE 2: SLM Architectures & SLM Models Matrix
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "SLM Architectures & Popular SLM Models", "SLM MODEL TAXONOMY")
    
    # Top Card: Architectures Summary
    top_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.2))
    top_box.fill.solid()
    top_box.fill.fore_color.rgb = CARD_BG
    top_box.line.color.rgb = CARD_BORDER
    tf = top_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = "Core SLM Architectures:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "1. Decoder-Only (Autoregressive - e.g. Phi-3, Llama 3.2)  |  2. Encoder-Only (Contextual - e.g. MobileBERT)  |  3. Encoder-Decoder (Seq2Seq - e.g. NLLB-200 600M, T5-Small)"
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_WHITE
    
    # Table for SLM Models
    rows, cols = 5, 7
    left, top_t, width, height = Inches(0.8), Inches(2.8), Inches(11.7), Inches(4.2)
    table_shape = slide2.shapes.add_table(rows, cols, left, top_t, width, height)
    table = table_shape.table
    
    # Column widths
    col_widths = [Inches(1.5), Inches(0.7), Inches(2.2), Inches(1.8), Inches(2.2), Inches(1.8), Inches(1.5)]
    for idx, w in enumerate(col_widths):
        table.columns[idx].width = w
        
    headers = ["Model Name", "Year", "Characteristics", "Companies Using", "Limitations", "Products Created", "Architecture"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 58, 138) # Dark blue header
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    models_data = [
        ("NLLB-200 Distilled (Used in Project)", "2022", "600M params, Seq2Seq multilingual translation across 200+ languages including Malayalam.", "Meta, Wikimedia Foundation, NLP Researchers", "Focused solely on translation; lacks multi-turn conversational reasoning.", "Facebook/Insta post translation, Wikipedia translation tool", "Encoder-Decoder"),
        ("Phi-3-mini", "2024", "3.8B params, state-of-the-art synthetic data training, 128k context window.", "Microsoft, Azure AI developers, Healthcare apps", "Higher memory than 600M models; slight hallucination on complex math.", "Microsoft Copilot local features, Azure AI Edge", "Decoder-only"),
        ("Llama-3.2 (1B/3B)", "2024", "1B & 3B lightweight params, Qualcomm/MediaTek mobile optimization, text+vision.", "Meta, Qualcomm, MediaTek, Edge AI developers", "Lower world knowledge capacity than Llama 70B.", "Meta AI mobile assistant, On-device edge features", "Decoder-only"),
        ("Gemma-2B", "2024", "2B params, built from Gemini tech, open weights, safety aligned.", "Google, Kaggle community, Android app devs", "Limited long-context multi-step reasoning capabilities.", "Google AI Edge, Android AICore features", "Decoder-only")
    ]
    
    for row_idx, data in enumerate(models_data, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if row_idx % 2 == 1 else RGBColor(15, 23, 42)
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(9.5)
            p.font.color.rgb = ACCENT_GREEN if col_idx == 0 and "Used" in text else TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 3: How LLM Works
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "How Large Language Model (LLM) Works", "LARGE LANGUAGE MODEL OVERVIEW")
    
    # Main Box
    llm_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.4))
    llm_box.fill.solid()
    llm_box.fill.fore_color.rgb = CARD_BG
    llm_box.line.color.rgb = CARD_BORDER
    
    tf = llm_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "Understanding Large Language Models (LLMs)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    sections = [
        ("1. Scale & Deep Architecture", "LLMs (e.g. GPT-4, Gemini 1.5, Claude 3.5) consist of tens to hundreds of billions of parameters (70B - 1 Trillion+). Built on stacked Transformer blocks with Multi-Head Self-Attention layers."),
        ("2. Massive Web-Scale Pre-Training", "Trained on trillions of tokens from web crawl datasets, books, code repositories, and scientific literature to learn deep world knowledge and grammar patterns."),
        ("3. Training Pipeline", "Pre-Training (Next-Token Prediction) → Supervised Fine-Tuning (SFT for instruction following) → Reinforcement Learning from Human Feedback (RLHF / DPO for safety & alignment)."),
        ("4. Inference & Generation Process", "Input Prompt → Tokenization (BPE/WordPiece) → High-dimensional Vector Embedding → Stacked Attention Layers → Softmax Probability Logits → Autoregressive Next-Token Sampling → Output Response."),
        ("5. Commercial Examples & Products", "ChatGPT (OpenAI), Gemini (Google), Claude (Anthropic), Llama 3.1 405B (Meta).")
    ]
    
    for title, desc in sections:
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(8)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 4: SLM vs LLM
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Comparative Analysis: SLM vs LLM", "MODEL ARCHITECTURE COMPARISON")
    
    # Table for SLM vs LLM
    rows, cols = 8, 4
    left, top_t, width, height = Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.4)
    table_shape = slide4.shapes.add_table(rows, cols, left, top_t, width, height)
    table = table_shape.table
    
    col_widths = [Inches(2.5), Inches(4.5), Inches(4.5), Inches(0.2)]
    table.columns[0].width = Inches(2.7)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(4.5)
    
    headers = ["Feature / Parameter", "Small Language Model (SLM)", "Large Language Model (LLM)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    comparison_data = [
        ("Parameter Size", "Under 10 Billion parameters (e.g. 600M, 1B, 3.8B)", "Tens to Hundreds of Billions (70B - 1Trillion+)"),
        ("Hardware Requirements", "Runs on edge devices, CPU, laptop, or low-cost single GPU", "Requires massive multi-GPU / TPU clusters (e.g. 8x H100s)"),
        ("Inference Latency & Speed", "Ultra-fast response time (milliseconds, real-time)", "Higher latency, queue dependent on cloud infrastructure"),
        ("Training & Operating Cost", "Low cost ($1k - $50k training), minimal operational overhead", "Extremely expensive ($1M - $100M+ per model run)"),
        ("Data Privacy & Security", "Complete local offline processing; zero data leakage risk", "Data sent via Cloud API endpoints; risk of telemetry leakage"),
        ("Domain Adaptability", "Highly agile; easily fine-tuned on custom domain datasets", "Generalist capability; complex & costly to fine-tune directly"),
        ("Best Use Cases", "On-device AI, Machine Translation (e.g. Malayalam), IoT, Privacy NLP", "Complex multi-step reasoning, general QA, coding assistants")
    ]
    
    for row_idx, (feat, slm_val, llm_val) in enumerate(comparison_data, start=1):
        # Feature column
        cell0 = table.cell(row_idx, 0)
        cell0.fill.solid()
        cell0.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p0 = cell0.text_frame.paragraphs[0]
        p0.text = feat
        p0.font.bold = True
        p0.font.size = Pt(11)
        p0.font.color.rgb = ACCENT_BLUE
        
        # SLM column
        cell1 = table.cell(row_idx, 1)
        cell1.fill.solid()
        cell1.fill.fore_color.rgb = RGBColor(15, 23, 42)
        p1 = cell1.text_frame.paragraphs[0]
        p1.text = slm_val
        p1.font.size = Pt(10.5)
        p1.font.color.rgb = ACCENT_GREEN
        
        # LLM column
        cell2 = table.cell(row_idx, 2)
        cell2.fill.solid()
        cell2.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p2 = cell2.text_frame.paragraphs[0]
        p2.text = llm_val
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 5: English Glossary to Malayalam Using SLM
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Glossary Translation: English → Malayalam (SLM Program)", "PRACTICAL IMPLEMENTATION & RUNTIME RESULTS")
    
    # Left Box: Program & Model Context
    box_l = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.0), Inches(5.6))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = CARD_BG
    box_l.line.color.rgb = CARD_BORDER
    tf = box_l.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
    
    p = tf.paragraphs[0]
    p.text = "English to Malayalam SLM Engine"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    info_points = [
        "Domain Context: Computer Science, AI & Natural Language Processing",
        "PDF Input: glossary_en.pdf (15 technical terms + definitions)",
        "SLM Model: facebook/nllb-200-distilled-600M (Meta)",
        "Target Language: Malayalam (mal_Mlym)",
        "Key Translation Examples:"
    ]
    for ip in info_points:
        p = tf.add_paragraph()
        p.text = "• " + ip
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(4)
        
    # Table of sample Malayalam translations inside left card
    samples = [
        ("Tokenization", "ടോക്കണൈസേഷൻ"),
        ("Small Language Model", "ചെറിയ ഭാഷാ മാതൃക (SLM)"),
        ("Attention Mechanism", "അറ്റൻഷൻ മെക്കാനിസം"),
        ("Knowledge Distillation", "നോളജ് ഡിസ്റ്റിലേഷൻ"),
        ("Quantization", "ക്വാണ്ടൈസേഷൻ")
    ]
    for en, ml in samples:
        p = tf.add_paragraph()
        p.text = f"  └ {en} → {ml}"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN

    # Right Side: Program Output Screenshot
    if os.path.exists(screenshot_path):
        slide5.shapes.add_picture(screenshot_path, Inches(6.0), Inches(1.4), Inches(6.5), Inches(5.6))
    else:
        # Placeholder box if screenshot not yet generated
        box_r = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(1.4), Inches(6.5), Inches(5.6))
        box_r.fill.solid()
        box_r.fill.fore_color.rgb = CARD_BG
        tf_r = box_r.text_frame
        p = tf_r.paragraphs[0]
        p.text = "[PROGRAM EXECUTION SCREENSHOT PLACEHOLDER]\nRun translate_glossary.py to render execution screenshot."
        p.font.color.rgb = ACCENT_GREEN

    # ----------------------------------------------------
    # SLIDE 6: GitHub Repository
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "GitHub Repository & Code Submission", "PROJECT REPOSITORY & DELIVERABLES")
    
    # Center GitHub Card
    gh_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.33), Inches(5.4))
    gh_card.fill.solid()
    gh_card.fill.fore_color.rgb = CARD_BG
    gh_card.line.color.rgb = CARD_BORDER
    
    tf = gh_card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "Project GitHub Repository Details"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    p_url_label = tf.add_paragraph()
    p_url_label.text = "GitHub Repository URL:"
    p_url_label.font.size = Pt(14)
    p_url_label.font.bold = True
    p_url_label.font.color.rgb = TEXT_WHITE
    p_url_label.space_before = Pt(10)
    
    p_url = tf.add_paragraph()
    p_url.text = "https://github.com/Awaaz-123/slm-english-malayalam-glossary.git"
    p_url.font.size = Pt(16)
    p_url.font.bold = True
    p_url.font.color.rgb = ACCENT_GREEN
    p_url.space_after = Pt(14)
    
    details = [
        ("Project Name", "English to Malayalam AI/NLP Glossary Translation using SLM"),
        ("Language & Frameworks", "Python 3.10+, Meta NLLB-200-distilled-600M, HuggingFace Transformers, PyTorch, PyPDF"),
        ("Key Deliverables", "1. translate_glossary.py (SLM Translation Script)\n2. CSE24562_Malayalam.pptx (6-Slide Presentation Deck)\n3. glossary_translation_results.json / .md (Translated Datasets)\n4. execution_screenshot.png (Program Execution Proof)"),
        ("Native Language", "Malayalam (mal_Mlym)"),
        ("Result", "Successfully parsed English PDF glossary and translated 15 complex technical AI/NLP definitions into Malayalam using a 600M parameter SLM.")
    ]
    
    for title, val in details:
        p_t = tf.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.bold = True
        p_t.font.size = Pt(12.5)
        p_t.font.color.rgb = ACCENT_BLUE
        
        p_v = tf.add_paragraph()
        p_v.text = f"  {val}"
        p_v.font.size = Pt(11.5)
        p_v.font.color.rgb = TEXT_WHITE
        p_v.space_after = Pt(4)

    # Save presentation
    prs.save(output_pptx)
    print(f"\n[Presentation Created] Successfully saved presentation to {output_pptx}")

if __name__ == "__main__":
    create_presentation()
